#!/usr/bin/env python3
"""
Aether_Access Build 2.0 - PowerPoint Presentation Generator
Generates a professional PowerPoint presentation for demo purposes
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_aetheraccess_presentation():
    """Create the complete Aether_Access presentation"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define colors
    NAVY = RGBColor(30, 58, 138)  # Dark blue
    GREEN = RGBColor(16, 185, 129)  # Success green
    ORANGE = RGBColor(245, 158, 11)  # Warning orange
    RED = RGBColor(239, 68, 68)  # Error red
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(31, 41, 55)  # Dark gray

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title background rectangle
    title_bg = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(2.5), Inches(10), Inches(2.5)
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = NAVY
    title_bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "BYTE_CONFIRM BUILD 2.0"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "Enterprise Access Control, Modernized"
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    tagline_frame = tagline_box.text_frame
    p = tagline_frame.paragraphs[0]
    p.text = "Surpassing Lenel OnGuard and Mercury Partners"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    # Slide 2: The Problem - Lenel OnGuard
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "LENEL ONGUARD: THE INDUSTRY STANDARD?"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = NAVY

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    problems = [
        ("Limited REST API", "Most functionality unavailable via API"),
        ("No Real-time Updates", "No WebSocket support, manual refresh required"),
        ("Dated User Interface", "Windows-style thick client, not mobile responsive"),
        ("Limited I/O Control via API", "Most operations GUI-only"),
        ("No Secure Channel Visibility", "Black box encryption, cannot diagnose issues"),
        ("Expensive", "$50K-$100K+ over 5 years")
    ]

    for problem, detail in problems:
        p = tf.add_paragraph()
        p.text = f"❌ {problem}"
        p.level = 0
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RED
        p.space_after = Pt(6)

        p = tf.add_paragraph()
        p.text = detail
        p.level = 1
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.space_after = Pt(12)

    # Slide 3: The Solution - Aether_Access
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "BYTE_CONFIRM BUILD 2.0: THE MODERN SOLUTION"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = NAVY

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    features = [
        "Complete REST API (30+ endpoints)",
        "Real-time WebSocket Updates",
        "Modern React Web Interface",
        "Full I/O Control via API",
        "Deep Secure Channel Visibility",
        "Comprehensive Health Monitoring",
        "Emergency Operations Built-in",
        "Integration Examples Included",
        "Production-Ready Monitoring Tools",
        "FREE - $0 Cost"
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = f"✅ {feature}"
        p.font.size = Pt(16)
        p.font.bold = (feature == features[-1])
        p.font.color.rgb = GREEN if feature == features[-1] else GRAY
        p.space_after = Pt(8)

    # Slide 4: Head-to-Head Comparison
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "HEAD-TO-HEAD COMPARISON"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER

    # Create comparison table
    rows = 13
    cols = 4
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5.5)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(2.5)

    # Header row
    headers = ["Feature", "Aether_Access 2.0", "Lenel OnGuard", "Mercury Partners"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data rows
    data = [
        ["Modern REST API", "✅ Complete", "⚠️ Limited", "⚠️ Varies"],
        ["WebSocket Real-time", "✅ Yes", "❌ No", "❌ Rare"],
        ["Modern Web UI", "✅ React", "⚠️ Dated", "⚠️ Varies"],
        ["I/O Control via API", "✅ Full", "⚠️ GUI only", "⚠️ Limited"],
        ["SC Deep Visibility", "✅ Deep", "❌ None", "⚠️ Basic"],
        ["Reader Health", "✅ 4 categories", "⚠️ Basic", "⚠️ Limited"],
        ["Auto Documentation", "✅ Swagger", "❌ No", "❌ Rare"],
        ["Open Architecture", "✅ Open", "❌ Proprietary", "⚠️ Proprietary"],
        ["Integration Examples", "✅ 3 languages", "⚠️ Limited", "⚠️ Minimal"],
        ["Emergency Operations", "✅ Built-in", "⚠️ Custom", "⚠️ Varies"],
        ["Mobile Responsive", "✅ Yes", "❌ No", "⚠️ Varies"],
        ["Cost (5-year)", "✅ $0", "❌ $50K-100K", "⚠️ $30K-60K"]
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Color code based on symbols
            if "✅" in cell_text:
                cell.text_frame.paragraphs[0].font.color.rgb = GREEN
            elif "❌" in cell_text:
                cell.text_frame.paragraphs[0].font.color.rgb = RED
            elif "⚠️" in cell_text:
                cell.text_frame.paragraphs[0].font.color.rgb = ORANGE

    # Slide 5: Modern Web Interface
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "MODERN WEB INTERFACE"
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.color.rgb = NAVY

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "React-Based Modern UI"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(12)

    features = [
        "✓ Real-time Dashboard with live metrics",
        "✓ Reader health monitoring (4 categories)",
        "✓ Live event feed via WebSocket",
        "✓ Responsive design - works on any device",
        "✓ No installation required - browser-based",
        "✓ Mobile responsive"
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(16)
        p.font.color.rgb = GREEN
        p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "\nCompare to:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = GRAY
    p.space_after = Pt(8)

    p = tf.add_paragraph()
    p.text = "❌ Lenel: Windows 95-style interface, thick client required"
    p.font.size = Pt(14)
    p.font.color.rgb = RED

    # Slide 6: Complete REST API
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "COMPLETE REST API WITH AUTO DOCUMENTATION"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = NAVY

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    p = tf.add_paragraph()
    p.text = "30+ Endpoints Covering:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(12)

    endpoints = [
        "Panel I/O Monitoring",
        "Reader Health (4 categories)",
        "Door Control (unlock, lock, lockdown)",
        "Output Control (activate, deactivate, pulse)",
        "Relay Control",
        "Control Macros",
        "Mass Operations (emergency)",
        "WebSocket Events"
    ]

    for endpoint in endpoints:
        p = tf.add_paragraph()
        p.text = f"• {endpoint}"
        p.font.size = Pt(14)
        p.font.color.rgb = GRAY
        p.space_after = Pt(6)

    p = tf.add_paragraph()
    p.text = "\n✅ Automatic Swagger/OpenAPI documentation"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = GREEN

    # Slide 7: Cost Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "5-YEAR TOTAL COST OF OWNERSHIP"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = NAVY

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    costs = [
        ("Lenel OnGuard:", "$50,000 - $100,000+", RED),
        ("Mercury + Partners:", "$30,000 - $60,000+", ORANGE),
        ("Aether_Access 2.0:", "$0", GREEN)
    ]

    for system, cost, color in costs:
        p = tf.add_paragraph()
        p.text = system
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = GRAY
        p.space_after = Pt(4)

        p = tf.add_paragraph()
        p.text = cost
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(20)

    p = tf.add_paragraph()
    p.text = "\nYOUR SAVINGS: $50,000 - $100,000"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_after = Pt(12)

    p = tf.add_paragraph()
    p.text = "What could you do with $50K-100K?"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = GRAY

    # Slide 8: Why Aether_Access Wins
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "WHY BYTE_CONFIRM BUILD 2.0 WINS"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = NAVY

    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()

    categories = [
        ("Superior Features:", [
            "30+ REST API endpoints vs limited API",
            "Real-time WebSocket vs no real-time",
            "Modern React UI vs dated Windows UI",
            "Full I/O control via API vs GUI-only"
        ]),
        ("Superior Visibility:", [
            "4-category health monitoring vs basic status",
            "Deep SC diagnostics vs black box",
            "Real-time event feed vs delayed updates"
        ]),
        ("Superior Integration:", [
            "Complete API documentation vs sparse docs",
            "3 language examples vs minimal examples",
            "WebSocket for real-time vs polling only"
        ]),
        ("Superior Cost:", [
            "$0 vs $50K-100K (Lenel)",
            "$0 vs $30K-60K (Mercury)"
        ])
    ]

    for category, items in categories:
        p = tf.add_paragraph()
        p.text = category
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY
        p.space_after = Pt(8)

        for item in items:
            p = tf.add_paragraph()
            p.text = f"✅ {item}"
            p.level = 1
            p.font.size = Pt(12)
            p.font.color.rgb = GREEN
            p.space_after = Pt(4)

    # Slide 9: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = "BYTE_CONFIRM BUILD 2.0"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = "Enterprise Access Control, Modernized"
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Key points
    points_box = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(6), Inches(2))
    points_frame = points_box.text_frame
    points_frame.word_wrap = True

    points = [
        "✅ Better Features",
        "✅ Better Usability",
        "✅ Better Integration",
        "✅ Better Visibility",
        "✅ Better Value"
    ]

    for point in points:
        p = points_frame.add_paragraph() if point != points[0] else points_frame.paragraphs[0]
        p.text = point
        p.font.size = Pt(18)
        p.font.color.rgb = GREEN
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

    # Final tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.6))
    tagline_frame = tagline_box.text_frame
    p = tagline_frame.paragraphs[0]
    p.text = "Aether_Access 2.0 >> Lenel OnGuard >> Mercury Partners"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_path = "/Users/mosley/Documents/AetherAccess_2.0_Presentation.pptx"
    prs.save(output_path)
    print(f"✅ PowerPoint presentation created: {output_path}")
    print(f"   - {len(prs.slides)} slides")
    print(f"   - Professional design with Aether_Access branding")
    print(f"   - Ready to present or customize further")

    return output_path

if __name__ == "__main__":
    create_aetheraccess_presentation()
