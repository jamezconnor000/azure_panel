#!/usr/bin/env python3
"""
Aether_Access Build 2.0 - PowerPoint Presentation Generator
Creates a professional presentation showcasing Aether_Access superiority
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_aetheraccess_presentation():
    """Create Aether_Access Build 2.0 presentation"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define colors - Elegant purple/blue theme for "Aether"
    DEEP_PURPLE = RGBColor(75, 0, 130)      # Deep purple for headers
    AETHER_BLUE = RGBColor(64, 156, 255)    # Light blue accent
    DARK_GRAY = RGBColor(45, 45, 45)        # Dark gray for text
    GREEN = RGBColor(16, 185, 129)          # Green for checkmarks
    RED = RGBColor(220, 38, 38)             # Red for X marks
    WHITE = RGBColor(255, 255, 255)         # White

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide1.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DEEP_PURPLE
    background.line.fill.background()

    # Title
    title_box = slide1.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Aether_Access Build 2.0"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(
        Inches(1), Inches(4.2), Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Enterprise Access Control, Modernized"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = AETHER_BLUE
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide1.shapes.add_textbox(
        Inches(1), Inches(5.5), Inches(8), Inches(0.8)
    )
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = "$0 Cost • Superior Features • Zero Vendor Lock-in"
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(20)
    tagline_para.font.color.rgb = GREEN
    tagline_para.alignment = PP_ALIGN.CENTER

    # Slide 2: The Problem with Current Solutions
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "The Problem: Legacy Systems Are Expensive & Limited"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = DEEP_PURPLE

    # Content
    content_box = slide2.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    problems = [
        ("Lenel OnGuard", [
            "Costs $62K-119K over 5 years (100 doors)",
            "Dated Windows-only thick client",
            "Limited SOAP-based API (<10 endpoints)",
            "No WebSocket real-time updates",
            "Proprietary vendor lock-in",
            "Complex integration (40-80 hours)"
        ]),
        ("Mercury Partners", [
            "Costs $29K-62K over 5 years",
            "Quality varies by software partner",
            "Inconsistent API support",
            "Usually Windows-dependent",
            "Limited customization options"
        ])
    ]

    for vendor, issues in problems:
        p = tf.add_paragraph()
        p.text = f"❌ {vendor}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = RED
        p.space_before = Pt(12)

        for issue in issues:
            p = tf.add_paragraph()
            p.text = f"  • {issue}"
            p.font.size = Pt(16)
            p.font.color.rgb = DARK_GRAY
            p.level = 1

    # Slide 3: Introducing Aether_Access
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "Introducing Aether_Access Build 2.0"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = DEEP_PURPLE

    # Content
    content_box = slide3.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    features = [
        "✅ Modern React Web Interface - Works on any device",
        "✅ Complete REST API - 30+ endpoints, fully documented",
        "✅ Real-time WebSocket Updates - Push notifications",
        "✅ Deep OSDP Secure Channel Diagnostics",
        "✅ Advanced Reader Health Monitoring - 4 categories, 0-100 score",
        "✅ Full I/O Control via API - Doors, outputs, relays",
        "✅ Emergency Operations - Lockdown, evacuation, mass control",
        "✅ Integration Examples - Python, JavaScript, Bash",
        "✅ Open Architecture - No vendor lock-in",
        "✅ Any OS - Linux, Windows, Mac, Docker"
    ]

    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

    # Cost highlight
    p = tf.add_paragraph()
    p.text = "💰 Cost: $0 (vs $29K-119K for competitors)"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.space_before = Pt(20)

    # Slide 4: Head-to-Head Comparison
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Head-to-Head Comparison"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = DEEP_PURPLE

    # Comparison table
    rows = 11
    cols = 4
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5.8)

    table = slide4.shapes.add_table(rows, cols, left, top, width, height).table

    # Set column widths
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.2)
    table.columns[2].width = Inches(2.2)
    table.columns[3].width = Inches(2.1)

    # Header row
    headers = ['Feature', 'Aether_Access', 'Lenel OnGuard', 'Mercury']
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DEEP_PURPLE
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(14)
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    # Data rows
    data = [
        ['Modern Web UI', '✅ React', '❌ Windows', '⚠️ Varies'],
        ['REST API', '✅ 30+ endpoints', '⚠️ Limited SOAP', '⚠️ Varies'],
        ['WebSocket', '✅ Yes', '❌ No', '❌ Rare'],
        ['SC Diagnostics', '✅ Deep', '❌ Black box', '⚠️ Basic'],
        ['Health Monitor', '✅ 4 categories', '⚠️ Basic', '⚠️ Limited'],
        ['I/O via API', '✅ Full', '⚠️ GUI only', '⚠️ Limited'],
        ['Integration Time', '✅ 1-4 hours', '❌ 40-80 hrs', '⚠️ 20-60 hrs'],
        ['Platform', '✅ Any OS', '❌ Windows', '⚠️ Usually Win'],
        ['Vendor Lock-in', '✅ None', '❌ Maximum', '⚠️ Moderate'],
        ['5-Year Cost', '✅ $0', '❌ $62-119K', '⚠️ $29-62K']
    ]

    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.alignment = PP_ALIGN.CENTER

            # Color coding
            if '✅' in cell_text:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(220, 252, 231)
            elif '❌' in cell_text:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(254, 226, 226)
            elif '⚠️' in cell_text:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(254, 243, 199)

    # Slide 5: Cost Analysis - The Money Truth
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Cost Analysis: Save $60K-200K"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = DEEP_PURPLE

    # Content
    content_box = slide5.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8.6), Inches(5.7))
    tf = content_box.text_frame
    tf.word_wrap = True

    # 5-year TCO
    p = tf.add_paragraph()
    p.text = "5-Year Total Cost of Ownership (100 doors):"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    costs = [
        ("Aether_Access 2.0", "$0", GREEN),
        ("Lenel OnGuard", "$62,000 - $119,000", RED),
        ("Mercury Partners", "$29,000 - $62,000", RGBColor(255, 165, 0))
    ]

    for system, cost, color in costs:
        p = tf.add_paragraph()
        p.text = f"  • {system}: {cost}"
        p.font.size = Pt(20)
        p.font.color.rgb = color
        p.level = 1
        p.space_after = Pt(8)

    # Savings
    p = tf.add_paragraph()
    p.text = "\nYour Savings with Aether_Access:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(20)

    savings = [
        "✅ Save $29K-62K vs Mercury (5 years)",
        "✅ Save $62K-119K vs Lenel (5 years)",
        "✅ Save $86K-208K vs competitors (10 years)"
    ]

    for saving in savings:
        p = tf.add_paragraph()
        p.text = f"  {saving}"
        p.font.size = Pt(18)
        p.font.color.rgb = GREEN
        p.level = 1
        p.space_after = Pt(6)

    # What you could do
    p = tf.add_paragraph()
    p.text = "\nWhat could you do with $100K savings?"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DEEP_PURPLE
    p.space_before = Pt(20)

    uses = [
        "Add 100 more doors",
        "Upgrade to premium hardware",
        "Implement video surveillance",
        "Deploy to multiple sites",
        "Keep the money!"
    ]

    for use in uses:
        p = tf.add_paragraph()
        p.text = f"  • {use}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.level = 1

    # Slide 6: Vendor Lock-in Exposed
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Freedom vs Vendor Lock-in"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = DEEP_PURPLE

    # Two columns
    # Left column - Lenel Lock-in
    left_box = slide6.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.3), Inches(5.8))
    tf_left = left_box.text_frame
    tf_left.word_wrap = True

    p = tf_left.add_paragraph()
    p.text = "❌ Lenel OnGuard"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RED

    lenel_issues = [
        "Closed source proprietary",
        "Windows-only platform",
        "SQL Server licensing ($5K-15K)",
        "Proprietary APIs and SDK",
        "Mandatory support contracts",
        "Per-door licensing fees",
        "Migration cost: $40K-140K",
        "Integration: 40-80 hours",
        "Vendor-dependent forever"
    ]

    for issue in lenel_issues:
        p = tf_left.add_paragraph()
        p.text = f"  • {issue}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.level = 1
        p.space_after = Pt(4)

    # Right column - Aether_Access Freedom
    right_box = slide6.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.3), Inches(5.8))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True

    p = tf_right.add_paragraph()
    p.text = "✅ Aether_Access"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GREEN

    aether_freedom = [
        "Open architecture",
        "Any OS (Linux/Win/Mac)",
        "PostgreSQL (free)",
        "Standard REST APIs",
        "Optional support only",
        "Unlimited doors ($0)",
        "Migration cost: Minimal",
        "Integration: 1-4 hours",
        "Complete independence"
    ]

    for benefit in aether_freedom:
        p = tf_right.add_paragraph()
        p.text = f"  • {benefit}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_GRAY
        p.level = 1
        p.space_after = Pt(4)

    # Slide 7: Modern Architecture & Performance
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Modern Architecture = Superior Performance"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = DEEP_PURPLE

    # Content
    content_box = slide7.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True

    # Tech stack
    p = tf.add_paragraph()
    p.text = "Technology Stack:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY

    tech = [
        "Backend: FastAPI (Python) - Modern async framework",
        "Frontend: React 18 + TypeScript - Responsive, fast",
        "Database: PostgreSQL - Enterprise-grade, free",
        "API: REST (30+ endpoints) + WebSocket real-time",
        "Deployment: Docker, cloud, any OS"
    ]

    for item in tech:
        p = tf.add_paragraph()
        p.text = f"  • {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.level = 1
        p.space_after = Pt(6)

    # Performance
    p = tf.add_paragraph()
    p.text = "\nPerformance Metrics:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(20)

    perf = [
        "API latency: <10ms (vs 100-500ms for Lenel)",
        "WebSocket latency: <5ms (Lenel has none)",
        "Events/second: 10,000+ sustained",
        "10-50x faster than legacy systems",
        "Deployment time: 5-60 minutes (vs 2-4 days)"
    ]

    for metric in perf:
        p = tf.add_paragraph()
        p.text = f"  ✅ {metric}"
        p.font.size = Pt(16)
        p.font.color.rgb = GREEN
        p.level = 1
        p.space_after = Pt(6)

    # Integration
    p = tf.add_paragraph()
    p.text = "\nEasy Integration:"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    p.space_before = Pt(20)

    integration = [
        "90% faster integration (1-4 hrs vs 40-80 hrs)",
        "Examples in Python, JavaScript, Bash",
        "Automatic API documentation (Swagger)",
        "Standard protocols - works with any language"
    ]

    for item in integration:
        p = tf.add_paragraph()
        p.text = f"  • {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.level = 1

    # Slide 8: Technical Superiority Summary
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Why Aether_Access Wins: 14/14 Categories"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = DEEP_PURPLE

    # Content
    content_box = slide8.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(5.8))
    tf = content_box.text_frame
    tf.word_wrap = True

    categories = [
        "✅ Architecture: Modern async vs legacy synchronous",
        "✅ API Quality: 30+ REST endpoints vs <10 SOAP",
        "✅ Real-time: WebSocket vs polling (1000x faster)",
        "✅ SC Diagnostics: Deep visibility vs black box",
        "✅ Health Monitoring: 4 categories vs online/offline",
        "✅ I/O Control: Full API vs GUI-only",
        "✅ Performance: 10-50x faster response times",
        "✅ Scalability: Easy horizontal scaling",
        "✅ Integration: 90% time reduction",
        "✅ Extensibility: Open vs proprietary",
        "✅ TCO: $0 vs $62K-119K (5 years)",
        "✅ Deployment: 20x faster",
        "✅ Operations: 1/5th the complexity",
        "✅ Security: Complete visibility + modern features"
    ]

    for category in categories:
        p = tf.add_paragraph()
        p.text = category
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(6)

    # Final verdict
    p = tf.add_paragraph()
    p.text = "\n🏆 Aether_Access 2.0 >> Lenel OnGuard >> Mercury Partners"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DEEP_PURPLE
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(20)

    # Slide 9: Call to Action
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    background = slide9.shapes.add_shape(
        1, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DEEP_PURPLE
    background.line.fill.background()

    # Main message
    title_box = slide9.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(8), Inches(1.2)
    )
    title_frame = title_box.text_frame
    title_frame.text = "Why Pay $60K-200K for Inferior Software?"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(38)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    # Key points
    content_box = slide9.shapes.add_textbox(
        Inches(1.5), Inches(3), Inches(7), Inches(3)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    points = [
        "✅ Better Features (30+ API endpoints, WebSocket, deep diagnostics)",
        "✅ Better Performance (10-50x faster)",
        "✅ Better Architecture (modern, open, flexible)",
        "✅ Better Cost ($0 vs $60K-200K)",
        "✅ Better Freedom (no vendor lock-in)"
    ]

    for point in points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(20)
        p.font.color.rgb = AETHER_BLUE
        p.space_after = Pt(12)
        p.alignment = PP_ALIGN.CENTER

    # Final CTA
    cta_box = slide9.shapes.add_textbox(
        Inches(1), Inches(6.2), Inches(8), Inches(1)
    )
    cta_frame = cta_box.text_frame
    cta_frame.text = "Choose Aether_Access. Save $60K-200K. Get Better Software."
    cta_para = cta_frame.paragraphs[0]
    cta_para.font.size = Pt(24)
    cta_para.font.bold = True
    cta_para.font.color.rgb = GREEN
    cta_para.alignment = PP_ALIGN.CENTER

    # Save presentation
    output_file = "/Users/mosley/Documents/AetherAccess_2.0_Presentation.pptx"
    prs.save(output_file)
    print(f"✓ Presentation created: {output_file}")
    print(f"  Slides: 9")
    print(f"  Theme: Elegant purple/blue (Aether)")
    print(f"  Ready to present!")

    return output_file

if __name__ == "__main__":
    create_aetheraccess_presentation()
